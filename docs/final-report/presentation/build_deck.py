# -*- coding: utf-8 -*-
"""
Rebuilds cybersim-defense-deck.pptx with a projector-optimized layout:
large readable type, brighter accent colors on a dark theme, consistent
grids, and corrected slide layouts. Content is preserved from the prior
deck (trimmed where needed so larger fonts fit cleanly).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (brightened for projector contrast on dark) ----
BG     = RGBColor(0x0A, 0x0E, 0x1A)   # slide background
CARD   = RGBColor(0x0E, 0x16, 0x28)   # card fill
BAND   = RGBColor(0x10, 0x19, 0x2E)   # architecture band fill
CYAN   = RGBColor(0x22, 0xF5, 0xFF)
RED    = RGBColor(0xFF, 0x5C, 0x5C)
BLUE   = RGBColor(0x5C, 0xC8, 0xFF)
GOLD   = RGBColor(0xE6, 0xC7, 0x66)
GREEN  = RGBColor(0x35, 0xE0, 0xA1)
TEXT   = RGBColor(0xF0, 0xF4, 0xFC)
MUTED  = RGBColor(0xAE, 0xB6, 0xCC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Segoe UI"

EMU_IN = 914400
prs = Presentation()
prs.slide_width  = Emu(int(13.333 * EMU_IN))
prs.slide_height = Emu(int(7.5 * EMU_IN))
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(s, l, t, w, h, fill=None, line=None, line_w=2.0, rounded=False, radius=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def tb(s, l, t, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """paras: list of dicts {text, size, color, bold, sa(space_after pt), ls(line spacing)}."""
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for i, pd in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if pd.get("sa") is not None:
            p.space_after = Pt(pd["sa"])
        p.space_before = Pt(0)
        if pd.get("ls"):
            p.line_spacing = pd["ls"]
        r = p.add_run()
        r.text = pd["text"]
        r.font.name = FONT
        r.font.size = Pt(pd["size"])
        r.font.bold = pd.get("bold", False)
        r.font.color.rgb = pd.get("color", TEXT)
    return box


def one(text, size, color, bold=False, ls=1.0, sa=0):
    return [{"text": text, "size": size, "color": color, "bold": bold, "ls": ls, "sa": sa}]


def bullets(items, size, color=TEXT, ls=1.05, sa=6):
    return [{"text": "•  " + it, "size": size, "color": color, "ls": ls, "sa": sa} for it in items]


def title_bar(s, text, accent, size=38, subtitle=None):
    rect(s, 0, 0, 13.333, 0.10, fill=accent)
    tb(s, 0.45, 0.18, 12.4, 0.75, one(text, size, accent, bold=True), anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        tb(s, 0.45, 0.95, 12.4, 0.5, one(subtitle, 22, MUTED), anchor=MSO_ANCHOR.MIDDLE)


def card(s, l, t, w, h, accent):
    return rect(s, l, t, w, h, fill=CARD, line=accent, line_w=2.0, rounded=True, radius=0.06)


def page_num(s, n, accent=MUTED):
    tb(s, 12.4, 7.05, 0.8, 0.35, one(str(n), 12, accent), align=PP_ALIGN.RIGHT)


# ============================================================ SLIDE 1 — TITLE
s = slide()
rect(s, 0, 0, 13.333, 0.10, fill=CYAN)
rect(s, 0, 7.40, 13.333, 0.10, fill=GOLD)
rect(s, 0, 0.10, 0.08, 7.30, fill=RED)
rect(s, 13.25, 0.10, 0.08, 7.30, fill=BLUE)
tb(s, 1.0, 1.55, 11.333, 1.6, one("CYBERSIM", 80, CYAN, bold=True), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tb(s, 1.0, 3.15, 11.333, 0.9, one("A Dual-Perspective Cybersecurity Training Platform", 30, TEXT),
   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 4.5, 4.25, 4.333, 0.035, fill=GOLD)
tb(s, 1.0, 4.5, 11.333, 0.5, one("University of Jordan   |   KASIT   |   Department of Computer Science", 20, MUTED),
   align=PP_ALIGN.CENTER)
tb(s, 1.0, 5.05, 11.333, 0.5, one("Graduation Project 2026", 20, MUTED), align=PP_ALIGN.CENTER)
tb(s, 1.0, 5.55, 11.333, 0.5, one("Supervised by: Dr. [Supervisor Name]", 20, MUTED), align=PP_ALIGN.CENTER)

# ====================================================== SLIDE 2 — THE DIVIDE
s = slide()
title_bar(s, "THE OFFENSIVE–DEFENSIVE DIVIDE", RED, size=36)
card(s, 0.4, 1.2, 5.85, 3.7, RED)
tb(s, 0.65, 1.35, 5.4, 0.55, one("Red Team  (Offensive)", 24, RED, bold=True))
tb(s, 0.65, 2.0, 5.4, 2.8, bullets([
    "Kali tools: nmap, sqlmap, Metasploit",
    "Isolated VM or lab environment",
    "Goal: find vulnerabilities",
    "No view of the defender's response",
], 19, TEXT, sa=12))
card(s, 7.1, 1.2, 5.85, 3.7, BLUE)
tb(s, 7.35, 1.35, 5.4, 0.55, one("Blue Team  (Defensive)", 24, BLUE, bold=True))
tb(s, 7.35, 2.0, 5.4, 2.8, bullets([
    "SIEM dashboards: Splunk, Elastic",
    "Separate course, separate lab",
    "Goal: detect anomalies",
    "No view of the attacker's methods",
], 19, TEXT, sa=12))
rect(s, 4.27, 5.15, 4.8, 0.78, fill=BG, line=GOLD, line_w=2.0, rounded=True, radius=0.5)
tb(s, 4.27, 5.15, 4.8, 0.78, one("←   Never connected   →", 22, GOLD, bold=True),
   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tb(s, 0.45, 6.35, 12.45, 0.8,
   one("Students graduate without grasping the causal link between an attack and its detection.", 21, GOLD),
   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 2)

# ================================================== SLIDE 3 — THE PLATFORM
s = slide()
title_bar(s, "THE CYBERSIM PLATFORM", CYAN, subtitle="One session. Two perspectives. Same live commands.")
feat = [
    ("\U0001F534 Red Terminal", "Real Kali container. Run nmap, sqlmap, gobuster against a live target.", RED),
    ("\U0001F535 Blue SIEM Feed", "Live stream — Suricata / Zeek / auditd alerts fire in real time as you attack.", BLUE),
    ("\U0001F916 AI Socratic Tutor", "Watches every command. Gives Socratic hints (L1→L3) without revealing answers.", CYAN),
    ("⚡ Instant Reset", "Docker sandbox resets in ~8 seconds. Fail safely, then try again.", GOLD),
    ("\U0001F4CA Debrief & Score", "Phase-by-phase score, flag evidence, and a learning-insights report.", GREEN),
    ("\U0001F504 One Timeline", "Red actions and Blue alerts share one timeline — see cause → effect instantly.", CYAN),
]
xs = [0.4, 4.7, 9.0]
ys = [1.65, 4.05]
for i, (h, b, ac) in enumerate(feat):
    l = xs[i % 3]
    t = ys[i // 3]
    card(s, l, t, 4.0, 2.15, ac)
    tb(s, l + 0.22, t + 0.18, 3.6, 0.55, one(h, 21, ac, bold=True))
    tb(s, l + 0.22, t + 0.78, 3.6, 1.25, one(b, 17, TEXT, ls=1.08))
page_num(s, 3)

# ============================================== SLIDE 4 — ARCHITECTURE
s = slide()
title_bar(s, "SYSTEM ARCHITECTURE", CYAN)
layers = [
    ("BROWSER", "React 18  ·  Vite  ·  Tailwind  ·  xterm.js  ·  Three.js 3D hero  ·  Zustand state"),
    ("FASTAPI BACKEND", "Python 3.11  ·  Async WebSocket  ·  Docker SDK  ·  JWT Auth  ·  Scenario Engine  ·  AI Monitor"),
    ("DATA LAYER", "PostgreSQL (sessions / scores)  ·  Redis (real-time pub/sub, WS)  ·  Elasticsearch + Filebeat"),
    ("SCENARIO CONTAINERS", "Kali  ·  NovaMed PHP/MariaDB  ·  Nexora Samba AD  ·  Orion GoPhish / victim-sim"),
]
arrows = ["↓  WebSocket / REST API",
          "↓  DB reads / writes, cache, events",
          "↓  Docker exec, SDK, network isolation"]
top = 1.45
band_h = 0.95
step = 1.27
for i, (lab, body) in enumerate(layers):
    t = top + i * step
    rect(s, 0.4, t, 12.55, band_h, fill=BAND, line=CYAN, line_w=1.75, rounded=True, radius=0.05)
    tb(s, 0.6, t, 3.1, band_h, one(lab, 17, CYAN, bold=True), anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 3.85, t, 8.85, band_h, one(body, 17, TEXT, ls=1.05), anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        tb(s, 3.85, t + band_h - 0.02, 8.85, 0.3, one(arrows[i], 14, MUTED), anchor=MSO_ANCHOR.MIDDLE)
tb(s, 0.4, 6.75, 12.55, 0.5,
   one("All scenario networks are internal-only (no internet). Kali → targets only. 9/9 containers hardened.",
       17, GOLD), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 4)

# ============================================== SLIDE 5 — RED TEAM WORKSPACE
def six_card_slide(title, accent, items, num):
    s = slide()
    title_bar(s, title, accent)
    xs = [0.4, 6.9]
    ys = [1.05, 2.95, 4.85]
    for i, (h, b) in enumerate(items):
        l = xs[i % 2]
        t = ys[i // 2]
        card(s, l, t, 6.0, 1.7, accent)
        tb(s, l + 0.22, t + 0.14, 5.6, 0.5, one(h, 20, accent, bold=True))
        tb(s, l + 0.22, t + 0.66, 5.6, 0.95, one(b, 17, TEXT, ls=1.05))
    page_num(s, num)
    return s

six_card_slide("RED TEAM WORKSPACE", RED, [
    ("Real Kali Terminal", "xterm.js + WebSocket → docker exec on cybersim-kali. Full nmap, sqlmap, gobuster, impacket toolkit."),
    ("PTES Phase Tracker", "Recon → Enum → Vuln ID → Exploit → Post-Exploit → Report. Phase gates enforce methodology."),
    ("ROE Scope Gate", "scope_enforcer.py blocks commands targeting out-of-scope IPs. Students can't escape the sandbox."),
    ("\U0001F6A9 Flag Discovery", "output_patterns.py scans PTY output; a nudge chip fires: 'Flag detected — submit for +N pts'."),
    ("AI Socratic Tutor", "L1 (–2) concept  ·  L2 (–5) hypothesis  ·  L3 (–10) scaffold. Never reveals the answer."),
    ("Collaborative Notes", "Markdown notebook tagged finding / evidence / ioc / remediation. Content feeds the debrief report."),
], 5)

# ============================================== SLIDE 6 — BLUE TEAM WORKSPACE
six_card_slide("BLUE TEAM WORKSPACE", BLUE, [
    ("Live SIEM Feed", "Redis pub/sub → WebSocket → SiemFeed.jsx. Every red command publishes ≥1 event within ms."),
    ("MITRE ATT&CK Tags", "Each event carries a technique (T1046, T1110.001…), severity (INFO→CRITICAL), and log source."),
    ("Alert Triage", "ForensicsWorkbench: classify each event true / false positive / investigating. Feeds the IR report."),
    ("Detection Bonuses", "Detect a scan within 5 minutes → bonus points. Blue scoring mirrors real SOC KPIs."),
    ("aria-live Stream", "Feed is aria-live='polite' — screen readers announce new events. WCAG 2.2 AA compliant."),
    ("Scenario Coverage", "SC-01: 27 detection rules  ·  SC-02: 22 (4768/4662/4670)  ·  SC-03: 18 + C2 callback."),
], 6)

# ============================================== SLIDE 7 — AI SAFETY MODEL
s = slide()
title_bar(s, "AI SAFETY MODEL  —  OWASP LLM TOP 10", GOLD, size=32)
rows = [
    ("L0 — ROE Scope Gate", "scope_enforcer.py blocks out-of-scope target IPs before the AI ever sees the command."),
    ("L1 — Prompt Injection Barrier", "Student input wrapped in <UNTRUSTED_DATA> tags — system prompt never treats it as instructions."),
    ("L2 — Output Post-Filter", "Regex strips payload syntax: ../ traversal, SQL tautologies, literal credentials, flag values."),
    ("L3 — Budget Enforcement", "Per-user token budget (LLM10). Rate-limited: 1 tutor call / 10s, 1 nudge / 60s. Never per-keystroke."),
    ("L4 — Output Validation", "validate_ai_output() checks each reply against scenario secrets → safe static Socratic fallback."),
    ("Regression Tests", "OWASP-LLM adversarial suite — injection, credential extraction, flag fishing. 340 tests, all green."),
]
top = 0.98
rh = 0.92
gap = 0.07
for i, (lab, body) in enumerate(rows):
    t = top + i * (rh + gap)
    rect(s, 0.3, t, 12.73, rh, fill=BAND, line=GOLD, line_w=1.6, rounded=True, radius=0.05)
    tb(s, 0.5, t, 3.4, rh, one(lab, 18, GOLD, bold=True), anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 3.95, t, 8.9, rh, one(body, 17, TEXT, ls=1.0), anchor=MSO_ANCHOR.MIDDLE)
tb(s, 0.3, 6.92, 12.73, 0.45,
   one("Design principle: the AI is a Socratic guide, never a solution engine — it knows the environment but teaches, never reveals.",
       16, MUTED), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 7)

# ============================================== SLIDE 8 — THREE SCENARIOS
s = slide()
title_bar(s, "THREE TRAINING SCENARIOS", CYAN)
cols = [
    (0.3, "SC-01", "NovaMed Healthcare", "Web Application Pentest", RED,
     ["LFI → /etc/passwd", "SQLi on the login form", "Unauth Redis (port 6379)", "IDOR on patient records", "Backup leak: .env.bak JWT"],
     ["ModSecurity WAF alerts", "Apache access-log anomalies", "MariaDB query monitoring"]),
    (4.65, "SC-02", "Nexora AD Domain", "Active Directory Compromise", GOLD,
     ["AS-REP Roasting (no preauth)", "GPP cPassword (SYSVOL)", "Kerberoasting (MSSQLSvc SPN)", "DCSync (secretsdump krbtgt)"],
     ["4768 / 4769 Kerberos events", "4662 LDAP enumeration", "4670 privilege change"]),
    (9.0, "SC-03", "Orion Logistics", "Phishing Campaign", BLUE,
     ["theHarvester OSINT", "GoPhish campaign (:3333)", "SWAKS SMTP delivery", "Victim simulator → C2"],
     ["SPF / DMARC probe detection", "Phishing click telemetry", "Reverse-shell handler events"]),
]
for (l, code, name, kind, ac, atk, blu) in cols:
    card(s, l, 1.6, 4.03, 5.65, ac)
    tb(s, l + 0.18, 1.72, 3.7, 0.55, one(code, 26, ac, bold=True))
    tb(s, l + 0.18, 2.28, 3.7, 0.45, one(name, 19, TEXT, bold=True))
    tb(s, l + 0.18, 2.72, 3.7, 0.38, one(kind, 15, MUTED))
    tb(s, l + 0.18, 3.18, 3.7, 0.34, one("Attack paths", 16, ac, bold=True))
    tb(s, l + 0.18, 3.52, 3.7, 2.0, bullets(atk, 15, TEXT, sa=4))
    tb(s, l + 0.18, 5.55, 3.7, 0.34, one("Blue detections", 16, BLUE, bold=True))
    tb(s, l + 0.18, 5.89, 3.7, 1.3, bullets(blu, 15, TEXT, sa=4))
page_num(s, 8)

# ============================================== SLIDE 9 — SECURITY MODEL
six_card_slide("SECURITY & ISOLATION MODEL", CYAN, [
    ("Network Isolation", "3 internal-only Docker bridge networks. Zero outbound internet — verified by verify-network-isolation.sh."),
    ("Container Hardening", "7/9 containers: cap_drop ALL + no-new-privileges. SC-02 Samba left fail-open (documented rationale)."),
    ("Auth & Rate Limiting", "JWT RS256. Auth 5/min  ·  AI tutor 1 / 10s  ·  flags 10/min. All limits enforced via Redis."),
    ("STRIDE Threat Model", "6 residual risks (R1–R6) documented. R1 docker.sock deferred; R3 hardening improved this sprint."),
    ("No Real Exploits in Source", "Attacks driven by YAML config + structured output patterns. No live malware, ransomware, or C2 in source."),
    ("Content-Security-Policy", "CSP report-only header in nginx blocks XSS. style-src 'unsafe-inline' for xterm.js (documented)."),
], 9)

# ============================================== SLIDE 10 — RESULTS (fixed grid)
s = slide()
title_bar(s, "RESULTS & METRICS", CYAN)
stats = [
    ("358", "Backend tests passing\n(316 unit + 42 integration)", CYAN),
    ("46", "Vitest frontend tests\n(component + hook coverage)", BLUE),
    ("22", "Architecture diagrams\n(rendered SVG + PNG)", GOLD),
    ("3", "Complete scenarios\nwith kill-chain evidence", GREEN),
    ("<3s", "AI tutor p50 latency\n(OpenRouter / DeepSeek)", BLUE),
    ("98/100", "Self-assessed completion\n(graduation-ready)", CYAN),
]
xs = [0.55, 4.67, 8.78]
ys = [1.35, 4.05]
cw, ch = 4.0, 2.45
for i, (num, cap, ac) in enumerate(stats):
    l = xs[i % 3]
    t = ys[i // 3]
    card(s, l, t, cw, ch, ac)
    tb(s, l, t + 0.22, cw, 1.2, one(num, 60, ac, bold=True), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cap_lines = cap.split("\n")
    tb(s, l, t + 1.42, cw, 0.95,
       [{"text": cap_lines[0], "size": 18, "color": TEXT, "bold": True, "ls": 1.05, "sa": 2},
        {"text": cap_lines[1], "size": 15, "color": MUTED, "ls": 1.05, "sa": 0}],
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
page_num(s, 10)

# ============================================== SLIDE 11 — LIVE DEMO
s = slide()
title_bar(s, "LIVE DEMO — SC-01 KILL CHAIN", CYAN)
steps = [
    ("1. Login", "admin / CyberSimAdmin!"),
    ("2. Select SC-01", "NovaMed Healthcare scenario"),
    ("3. ROE Brief", "Acknowledge rules of engagement"),
    ("4. Recon", "nmap -sV 172.20.1.20"),
    ("5. SIEM Alert", "Blue: Port scan detected (LOW)"),
    ("6. AI Hint", "Ask: 'What should I enumerate next?'"),
    ("7. LFI", "curl .../records/?file=../../etc/passwd"),
    ("8. Flag Nudge", "\U0001F6A9 chip: 'root:x:0:0' detected"),
    ("9. Submit Flag", "FLAG-SC01-1 captured  ·  +15 pts"),
    ("10. Debrief", "Score: 100 → penalties / bonuses → final"),
]
xs = [0.3, 6.83]
ys = [1.0, 2.2, 3.4, 4.6, 5.8]
for i, (lab, detail) in enumerate(steps):
    l = xs[i % 2]
    t = ys[i // 2]
    rect(s, l, t, 6.2, 1.05, fill=BAND, line=CYAN, line_w=1.6, rounded=True, radius=0.06)
    tb(s, l + 0.18, t, 2.05, 1.05, one(lab, 18, CYAN, bold=True), anchor=MSO_ANCHOR.MIDDLE)
    tb(s, l + 2.25, t, 3.85, 1.05, one(detail, 16, TEXT, ls=1.0), anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 11)

# ============================================== SLIDE 12 — THANK YOU
s = slide()
rect(s, 0, 0, 13.333, 0.10, fill=CYAN)
rect(s, 0, 7.40, 13.333, 0.10, fill=GOLD)
rect(s, 0, 0.10, 0.08, 7.30, fill=RED)
rect(s, 13.25, 0.10, 0.08, 7.30, fill=BLUE)
tb(s, 1.0, 1.7, 11.333, 1.3, one("THANK YOU", 76, CYAN, bold=True), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tb(s, 1.0, 3.05, 11.333, 0.8, one("Questions?", 34, GOLD), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 3.5, 4.15, 6.333, 0.035, fill=GOLD)
info = [("GitHub:", "github.com/VinsmokeD/JUTerminal1"),
        ("Stack:", "React + FastAPI + Docker + OpenRouter/DeepSeek"),
        ("Tests:", "358 passing  |  46 Vitest  |  3 scenarios  |  98/100")]
iy = 4.5
for k, v in info:
    tb(s, 2.9, iy, 1.7, 0.5, one(k, 18, GOLD, bold=True), align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 4.75, iy, 6.0, 0.5, one(v, 18, TEXT), anchor=MSO_ANCHOR.MIDDLE)
    iy += 0.62

prs.save("cybersim-defense-deck.pptx")
print("Saved cybersim-defense-deck.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
