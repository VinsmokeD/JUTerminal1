# -*- coding: utf-8 -*-
"""Faithful-ish PNG preview of the deck using real Segoe UI metrics + word wrap.
Also flags any paragraph whose wrapped text exceeds its textbox height (overflow)."""
import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont

SCALE = 96  # px per inch
REG = "C:/Windows/Fonts/segoeui.ttf"
BLD = "C:/Windows/Fonts/segoeuib.ttf"
EMU = 914400
_fc = {}

def font(sz, bold):
    key = (int(sz), bold)
    if key not in _fc:
        _fc[key] = ImageFont.truetype(BLD if bold else REG, int(sz * SCALE / 72))
    return _fc[key]

def px(v):
    return int(v / EMU * SCALE)

def rgb(c):
    if not isinstance(c, int):
        c = int(str(c), 16)
    return (c >> 16 & 255, c >> 8 & 255, c & 255)

def wrap(draw, text, fnt, maxw):
    words = text.split(" ")
    lines, cur = [], ""
    for w in words:
        trial = (cur + " " + w).strip()
        if draw.textlength(trial, font=fnt) <= maxw or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines

def shape_fill(sh):
    try:
        if sh.fill.type == 1:
            return rgb(int(sh.fill.fore_color.rgb))
    except Exception:
        pass
    return None

def shape_line(sh):
    try:
        return rgb(int(sh.line.color.rgb)), max(1, px(sh.line.width))
    except Exception:
        return None, 0

p = Presentation("cybersim-defense-deck.pptx")
W = px(p.slide_width); H = px(p.slide_height)
os.makedirs("_preview", exist_ok=True)
overflows = []

for idx, s in enumerate(p.slides):
    img = Image.new("RGB", (W, H), rgb(0x0A0E1A))
    try:
        if s.background.fill.type == 1:
            img = Image.new("RGB", (W, H), rgb(int(s.background.fill.fore_color.rgb)))
    except Exception:
        pass
    d = ImageDraw.Draw(img)
    for sh in s.shapes:
        l, t, w, h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
        fill = shape_fill(sh)
        lc, lw = shape_line(sh)
        if fill is not None or lc is not None:
            d.rectangle([l, t, l + w, h + t], fill=fill, outline=lc, width=lw or 1)
        if sh.has_text_frame and sh.text_frame.text.strip():
            tf = sh.text_frame
            # build wrapped lines with heights
            blocks = []
            for para in tf.paragraphs:
                txt = "".join(r.text for r in para.runs)
                if not txt.strip():
                    continue
                r0 = para.runs[0]
                sz = r0.font.size.pt if r0.font.size else 18
                bold = bool(r0.font.bold)
                col = rgb(r0.font.color.rgb) if (r0.font.color and r0.font.color.type) else rgb(0xF0F4FC)
                al = para.alignment
                fnt = font(sz, bold)
                lh = int(sz * SCALE / 72 * (para.line_spacing or 1.08))
                sa = int((para.space_after.pt if para.space_after else 0) * SCALE / 72)
                wl = wrap(d, txt, fnt, w - px(40000))
                for j, ln in enumerate(wl):
                    extra = sa if j == len(wl) - 1 else 0
                    blocks.append((ln, fnt, col, al, lh + extra))
            total = sum(b[4] for b in blocks)
            anchor = tf.vertical_anchor
            y = t + px(12000)
            if anchor == MSO_ANCHOR.MIDDLE:
                y = t + max(0, (h - total) // 2)
            elif anchor == MSO_ANCHOR.BOTTOM:
                y = t + h - total
            if total > h + 4:
                overflows.append((idx + 1, sh.name, blocks[0][0][:32], round(total/SCALE,2), round(h/SCALE,2)))
            for ln, fnt, col, al, lh in blocks:
                tw = d.textlength(ln, font=fnt)
                x = l + px(20000)
                if al == PP_ALIGN.CENTER:
                    x = l + (w - tw) // 2
                elif al == PP_ALIGN.RIGHT:
                    x = l + w - tw - px(20000)
                d.text((x, y), ln, font=fnt, fill=col)
                y += lh
    img.save(f"_preview/slide_{idx+1:02d}.png")

print("Rendered", len(p.slides), "slides to _preview/")
print("\nTEXT OVERFLOW (wrapped height > box):")
if overflows:
    for o in overflows:
        print(f"  slide {o[0]} '{o[1]}' [{o[2]}...] needs {o[3]}in, box {o[4]}in")
else:
    print("  NONE")
